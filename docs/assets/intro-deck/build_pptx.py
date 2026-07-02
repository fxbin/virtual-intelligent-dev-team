#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
virtual-intelligent-dev-team 对外宣讲 PPTX 生成脚本。

使用 python-pptx 原生 shape 绘制 10 页深色风格幻灯片，
不依赖任何外部图片，所有视觉元素由 shape/textbox 构造。

@author: fxbin
@date: 2026-07-02
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================
# 颜色常量（machine-hallucination 配色，取自 ppt-creation-forge 色彩搭配库）
# paper=#0A0A1A 数据深空 / ink=#9B4DCA 数据紫 / accent=#00CED1 数据青
# muted=#4B0082 深紫 / secondary=#FF1493 数据粉
# PPT 背景提一档到 #12122A 防投影衰减；ink 紫 4.0:1 未达 WCAG AA，降级为装饰色
# 正文用 #c4b5fd 浅紫（10.6:1 AAA）；KPI/链接用 #00CED1 数据青（签名色）
# ============================================================
BG_DARK = RGBColor(0x12, 0x12, 0x2A)        # paper 提一档，防投影衰减
PURPLE = RGBColor(0x9B, 0x4D, 0xCA)         # ink 数据紫（装饰色，不做正文）
BLUE = RGBColor(0x9B, 0x4D, 0xCA)           # 合并到数据紫
CYAN = RGBColor(0x00, 0xCE, 0xD1)           # accent 数据青（签名色）
GREEN = RGBColor(0x00, 0xCE, 0xD1)          # 合并到数据青
AMBER = RGBColor(0xFF, 0x14, 0x93)          # secondary 数据粉（点睛色）
ROSE = RGBColor(0xFF, 0x14, 0x93)           # 合并到数据粉
INDIGO = RGBColor(0x9B, 0x4D, 0xCA)         # 合并到数据紫

WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # 标题纯白
TEXT_SECONDARY = RGBColor(0xC4, 0xB5, 0xFD) # 正文浅紫（10.6:1 AAA）
TEXT_MUTED = RGBColor(0x8B, 0x6F, 0xC0)     # 弱文字暗紫

CARD_BG = RGBColor(0x13, 0x10, 0x2A)        # 卡片背景微紫提亮
CARD_BORDER = RGBColor(0x2A, 0x1F, 0x4D)    # 边框低饱和深紫

# 字体（中文优先 PingFang SC，回退 Microsoft YaHei）
FONT_CN = "PingFang SC"
FONT_EN = "Helvetica"

# 画布尺寸（16:9）
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TOTAL_PAGES = 10


# ============================================================
# Helper 函数
# ============================================================
def set_solid_fill(shape, color):
    """为 shape 的 fill 设置纯色。"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_fill(shape):
    """关闭 shape 的填充。"""
    shape.fill.background()


def set_line(shape, color, width_pt=1.0):
    """为 shape 的边框设置颜色和宽度。"""
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)


def set_no_line(shape):
    """去除 shape 的边框。"""
    shape.line.fill.background()


def add_bg(slide, color=BG_DARK):
    """在 slide 上铺一层纯色背景矩形（覆盖整个画布）。"""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    set_solid_fill(rect, color)
    set_no_line(rect)
    # 移到最底层
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def _set_run(run, text, size, color, bold=False, font_name=FONT_CN):
    """配置一个 run 的文本与样式。"""
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=WHITE,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name=FONT_CN,
):
    """
    添加单段文本框。

    参数:
        slide: 目标幻灯片
        left/top/width/height: 位置与尺寸（EMU）
        text: 文本内容
        size: 字号 pt
        color: 字色
        bold: 是否加粗
        align: 对齐方式
        anchor: 垂直对齐
        font_name: 字体
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size, color, bold, font_name)
    return tb


def add_multi_run_text(
    slide,
    left,
    top,
    width,
    height,
    runs,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=None,
):
    """
    添加多 run 段落（用于一行内有多种颜色/字号的文本）。

    runs 是 list of dict: {text, size, color, bold, font_name}
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for spec in runs:
        _set_run(
            p.add_run(),
            spec["text"],
            spec.get("size", 18),
            spec.get("color", WHITE),
            spec.get("bold", False),
            spec.get("font_name", FONT_CN),
        )
    return tb


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    lines,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=None,
):
    """
    添加多段落文本框。

    lines: list of list(run_spec) — 每个元素是一个段落，由若干 run 组成。
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for idx, runs in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for spec in runs:
            _set_run(
                p.add_run(),
                spec["text"],
                spec.get("size", 18),
                spec.get("color", WHITE),
                spec.get("bold", False),
                spec.get("font_name", FONT_CN),
            )
    return tb


def add_card(
    slide,
    left,
    top,
    width,
    height,
    fill_color=CARD_BG,
    border_color=CARD_BORDER,
    border_width=1.0,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    """添加卡片背景（默认圆角矩形）。"""
    card = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_solid_fill(card, fill_color)
    set_line(card, border_color, border_width)
    # 调整圆角半径（避免过圆）
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            card.adjustments[0] = 0.06
        except Exception:
            pass
    return card


def add_eyebrow(slide, text, page_num):
    """在左上角添加页眉小标签，右下角添加页码。"""
    # eyebrow
    add_text(
        slide,
        Inches(0.6),
        Inches(0.4),
        Inches(6),
        Inches(0.3),
        text,
        size=12,
        color=CYAN,
        bold=True,
    )
    # 页码
    add_text(
        slide,
        Inches(11.6),
        Inches(7.0),
        Inches(1.4),
        Inches(0.35),
        f"{page_num:02d} / {TOTAL_PAGES}",
        size=11,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_badge(slide, left, top, text, bg_color, text_color=WHITE, width=1.6, height=0.34):
    """添加一个小徽章（圆角矩形 + 文本）。"""
    bw = Inches(width)
    bh = Inches(height)
    badge = add_card(
        slide, left, top, bw, bh,
        fill_color=bg_color,
        border_color=bg_color,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    try:
        badge.adjustments[0] = 0.5
    except Exception:
        pass
    tf = badge.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), text, 10, text_color, bold=True, font_name=FONT_EN)
    return badge


def add_section_title(slide, top, runs, size=40, space_before=0):
    """添加主标题（支持多 run 上色）。"""
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for spec in runs:
        _set_run(
            p.add_run(),
            spec["text"],
            spec.get("size", size),
            spec.get("color", WHITE),
            spec.get("bold", True),
            spec.get("font_name", FONT_CN),
        )
    return tb


def add_subtitle(slide, top, text, color=TEXT_SECONDARY, size=18, width=12.1):
    """添加副标题。"""
    return add_text(
        slide,
        Inches(0.6),
        top,
        Inches(width),
        Inches(0.4),
        text,
        size=size,
        color=color,
    )


def add_card_text(card, lines, anchor=MSO_ANCHOR.TOP, padding=0.15):
    """
    在已存在的卡片 shape 中写入多段落文本。

    lines: list of list(run_spec)
    """
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(padding)
    tf.margin_right = Inches(padding)
    tf.margin_top = Inches(padding)
    tf.margin_bottom = Inches(padding)
    tf.vertical_anchor = anchor
    for idx, runs in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        for spec in runs:
            _set_run(
                p.add_run(),
                spec["text"],
                spec.get("size", 14),
                spec.get("color", WHITE),
                spec.get("bold", False),
                spec.get("font_name", FONT_CN),
            )
    return card


# ============================================================
# 各 Slide 内容构造
# ============================================================
def build_slide_1(prs):
    """封面。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # 右上角 4 个徽章
    badge_y = Inches(0.4)
    badge_specs = [
        ("v5.7.6", PURPLE),
        ("MIT License", GREEN),
        ("Production Ready", AMBER),
        ("Router Archetype", CYAN),
    ]
    badge_x = Inches(13.333) - Inches(0.2)
    # 从右向左排
    for text, color in badge_specs:
        w = 1.7 if len(text) > 10 else 1.4
        badge_x = badge_x - Inches(w + 0.1)
        add_badge(slide, badge_x, badge_y, text, color, width=w, height=0.32)

    # 主标题
    add_text(
        slide,
        Inches(0.6),
        Inches(2.2),
        Inches(12.1),
        Inches(1.4),
        "Virtual Intelligent Dev Team",
        size=60,
        color=WHITE,
        bold=True,
        font_name=FONT_EN,
    )

    # 副标题
    add_text(
        slide,
        Inches(0.6),
        Inches(3.5),
        Inches(12.1),
        Inches(0.6),
        "面向复杂软件工作的闭环协调层",
        size=28,
        color=TEXT_SECONDARY,
    )

    # 描述
    add_text(
        slide,
        Inches(0.6),
        Inches(4.25),
        Inches(12.1),
        Inches(0.7),
        "把专家路由 · 计划 · 执行 · 迭代 · Beta · Release · Feedback 收拢成一个可持续迭代的闭环工作流",
        size=16,
        color=TEXT_MUTED,
    )

    # KPI 数字横排（视觉锚点）
    kpi_y = Inches(5.2)
    kpi_data = [
        ("8", "Specialists"),
        ("7", "Closures"),
        ("13", "Languages"),
        ("7", "Workflows"),
    ]
    kpi_total_w = 12.13
    kpi_gap = 1.5
    kpi_n = len(kpi_data)
    kpi_w = (kpi_total_w - kpi_gap * (kpi_n - 1)) / kpi_n
    kpi_start_x = 0.6
    for idx, (num, label) in enumerate(kpi_data):
        x = Inches(kpi_start_x + (kpi_w + kpi_gap) * idx)
        add_text(
            slide,
            x,
            kpi_y,
            Inches(kpi_w),
            Inches(0.8),
            num,
            size=48,
            color=PURPLE,
            bold=True,
            font_name=FONT_EN,
        )
        add_text(
            slide,
            x,
            kpi_y + Inches(0.85),
            Inches(kpi_w),
            Inches(0.35),
            label.upper(),
            size=11,
            color=TEXT_MUTED,
            font_name=FONT_EN,
        )

    # 底部元信息
    add_text(
        slide,
        Inches(0.6),
        Inches(6.8),
        Inches(12.1),
        Inches(0.35),
        "github.com/fxbin/virtual-intelligent-dev-team  ·  v5.7.6  ·  MIT License",
        size=12,
        color=TEXT_MUTED,
    )

    # 页码（封面也加）
    add_text(
        slide,
        Inches(11.6),
        Inches(7.0),
        Inches(1.4),
        Inches(0.35),
        f"01 / {TOTAL_PAGES}",
        size=11,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def build_slide_2(prs):
    """问题痛点。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "PROBLEM · 痛点", 2)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "为什么需要", "size": 36, "color": WHITE, "bold": True},
            {"text": "闭环协调层", "size": 36, "color": PURPLE, "bold": True},
            {"text": "?", "size": 36, "color": WHITE, "bold": True},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "普通多专家提示词在复杂软件任务中有三个结构性缺口",
    )

    # 3 个并排卡片（玫红边框）
    card_y = Inches(2.6)
    card_h = Inches(4.2)
    card_w = Inches(3.85)
    gap = Inches(0.25)
    start_x = Inches(0.6)

    items = [
        ("01", "角色堆砌 ≠ 协同", "平铺多个角色视角，没有判断谁主负责、谁协同、是否需要治理。复杂任务里'谁来 lead'比'有几个专家'更重要。"),
        ("02", "单轮回答 ≠ 闭环", "只给建议，没有执行路径、恢复锚点和下一步。真实软件交付需要跨多轮、跨多天、可恢复的状态管理。"),
        ("03", "发布前判断 ≠ 全链路", "只覆盖'能不能发'，缺少 beta、release gate、completion evidence、post-release feedback 的连续治理链路。"),
    ]

    for idx, (num, title, desc) in enumerate(items):
        x = start_x + (card_w + gap) * idx
        card = add_card(
            slide, x, card_y, card_w, card_h,
            fill_color=CARD_BG,
            border_color=ROSE,
            border_width=1.25,
        )
        add_card_text(
            card,
            [
                [{"text": num, "size": 36, "color": ROSE, "bold": True, "font_name": FONT_EN}],
                [{"text": title, "size": 18, "color": WHITE, "bold": True}],
                [{"text": desc, "size": 13, "color": TEXT_SECONDARY}],
            ],
            padding=0.25,
        )


def build_slide_3(prs):
    """定位。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "POSITIONING · 定位", 3)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "不是角色扮演器，是", "size": 36, "color": WHITE, "bold": True},
            {"text": "闭环协调层", "size": 36, "color": PURPLE, "bold": True},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "从'换几个角色回答'升级到'打通复杂软件工作的关键链条'",
    )

    # 左侧（弱化，虚线边框）
    left_x = Inches(0.7)
    box_y = Inches(2.7)
    box_w = Inches(5.0)
    box_h = Inches(3.2)

    left_card = add_card(
        slide, left_x, box_y, box_w, box_h,
        fill_color=CARD_BG,
        border_color=TEXT_MUTED,
        border_width=1.0,
    )
    # 用 dashed 线模拟
    ln = left_card.line
    ln.color.rgb = TEXT_MUTED
    ln.width = Pt(1.0)
    lnElem = left_card.line._get_or_add_ln()
    # 设置 dash 样式
    existing = lnElem.find(qn('a:prstDash'))
    if existing is None:
        existing = etree.SubElement(lnElem, qn('a:prstDash'))
    existing.set('val', 'dash')

    add_card_text(
        left_card,
        [
            [{"text": "✕  ", "size": 28, "color": TEXT_MUTED, "bold": True, "font_name": FONT_EN},
             {"text": "Ordinary", "size": 22, "color": TEXT_MUTED, "bold": True, "font_name": FONT_EN}],
            [{"text": "多专家提示词", "size": 16, "color": TEXT_SECONDARY}],
            [{"text": "换几个角色视角回答问题", "size": 13, "color": TEXT_MUTED}],
            [{"text": "单轮建议 · 无执行路径 · 无闭环", "size": 13, "color": TEXT_MUTED}],
        ],
        padding=0.3,
    )

    # 中间箭头
    add_text(
        slide,
        Inches(5.85),
        Inches(3.9),
        Inches(1.6),
        Inches(0.8),
        "→",
        size=54,
        color=PURPLE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
    )

    # 右侧（强调，紫色边框）
    right_x = Inches(7.6)
    right_card = add_card(
        slide, right_x, box_y, box_w, box_h,
        fill_color=CARD_BG,
        border_color=PURPLE,
        border_width=2.0,
    )
    add_card_text(
        right_card,
        [
            [{"text": "✦  ", "size": 28, "color": PURPLE, "bold": True, "font_name": FONT_EN},
             {"text": "Virtual Dev Team", "size": 22, "color": WHITE, "bold": True, "font_name": FONT_EN}],
            [{"text": "闭环协调层", "size": 16, "color": PURPLE, "bold": True}],
            [{"text": "路由 + 计划 + 执行 + 迭代 + Beta", "size": 13, "color": TEXT_SECONDARY}],
            [{"text": "Release + Feedback + 状态恢复", "size": 13, "color": TEXT_SECONDARY}],
        ],
        padding=0.3,
    )

    # 底部提示框
    tip = add_card(
        slide, Inches(0.7), Inches(6.2), Inches(11.95), Inches(0.8),
        fill_color=CARD_BG,
        border_color=CARD_BORDER,
    )
    add_card_text(
        tip,
        [[
            {"text": "一句话理解：", "size": 14, "color": AMBER, "bold": True},
            {"text": "适合接手'单个专家已经不够、单轮回答也不够'的复杂软件任务。", "size": 14, "color": WHITE},
        ]],
        anchor=MSO_ANCHOR.MIDDLE,
        padding=0.2,
    )


def build_slide_4(prs):
    """七层闭环。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "ARCHITECTURE · 七层闭环", 4)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "Seven ", "size": 36, "color": WHITE, "bold": True, "font_name": FONT_EN},
            {"text": "Closure Layers", "size": 36, "color": PURPLE, "bold": True, "font_name": FONT_EN},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "每层都有明确触发条件、核心机制和结构化产物",
    )

    # 顶部 7 个小卡片横排
    layers = [
        ("L1", "规划 Planning", PURPLE),
        ("L2", "路由 Routing", BLUE),
        ("L3", "交付 Delivery", CYAN),
        ("L4", "迭代 Iteration", GREEN),
        ("L5", "发布 Release", AMBER),
        ("L6", "演练 Drill", ROSE),
        ("L7", "团队引擎 Team Engine Lite", INDIGO),
    ]
    card_y = Inches(2.55)
    card_h = Inches(1.05)
    total_w = 12.13
    gap = 0.08
    n = len(layers)
    card_w = (total_w - gap * (n - 1)) / n
    start_x = 0.6

    for idx, (label, name, color) in enumerate(layers):
        x = Inches(start_x + (card_w + gap) * idx)
        w = Inches(card_w)
        card = add_card(
            slide, x, card_y, w, card_h,
            fill_color=CARD_BG,
            border_color=CARD_BORDER,
        )
        # 顶部色条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, card_y, w, Inches(0.08)
        )
        set_solid_fill(top_bar, color)
        set_no_line(top_bar)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.18)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), label, 18, color, bold=True, font_name=FONT_EN)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _set_run(p2.add_run(), name, 11, WHITE, bold=False)

    # 下方 3 个说明卡片
    desc_y = Inches(4.0)
    desc_h = Inches(2.6)
    desc_w = Inches(3.85)
    desc_gap = Inches(0.25)
    descs = [
        ("L1-L2 · 入口层", PURPLE, "规划与路由：判断任务形态，选择主负责人和最小 workflow bundle。模糊想法先做意图确认。"),
        ("L3-L4 · 执行层", CYAN, "交付与迭代：小切片保留 brief + context + status；优化循环有 baseline、round memory 和回滚决策。"),
        ("L5-L7 · 治理层", AMBER, "发布、演练、团队引擎：ship/hold 门禁 + completion evidence；Worker/Verifier 分离验收。"),
    ]
    for idx, (title, color, body) in enumerate(descs):
        x = Inches(0.6) + (desc_w + desc_gap) * idx
        card = add_card(
            slide, x, desc_y, desc_w, desc_h,
            fill_color=CARD_BG,
            border_color=color,
            border_width=1.25,
        )
        add_card_text(
            card,
            [
                [{"text": title, "size": 16, "color": color, "bold": True}],
                [{"text": body, "size": 13, "color": TEXT_SECONDARY}],
            ],
            padding=0.25,
        )


def build_slide_5(prs):
    """8 个 Agent。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "AGENT ROSTER · 专家角色", 5)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "8 ", "size": 36, "color": WHITE, "bold": True, "font_name": FONT_EN},
            {"text": "Specialist Agents", "size": 36, "color": PURPLE, "bold": True, "font_name": FONT_EN},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "覆盖从代码到产品、从架构到发布的完整软件交付链路",
    )

    # 4×2 网格
    agents = [
        ("JV", "Java Virtuoso", "JVM · Spring · Gradle", ROSE),
        ("SA", "Sentinel Architect", "架构风险 · 高危治理", PURPLE),
        ("TT", "Technical Trinity", "通用后端 · 小切片", GREEN),
        ("CA", "Code Audit Council", "安全审计 · 代码审查", AMBER),
        ("GW", "Git Workflow Guardian", "分支策略 · 回滚保护", CYAN),
        ("PA", "World-Class Product Architect", "PRD · 用户研究 · 路线图", BLUE),
        ("DP", "Data Pipeline Guardian", "Kafka · Flink · 实时管道", ROSE),
        ("AC", "API Contract Sentinel", "契约 · 兼容性 · 版本", INDIGO),
    ]

    grid_x = 0.6
    grid_y = 2.6
    cols = 4
    rows = 2
    cell_w = 2.95
    cell_h = 1.65
    gap_x = 0.13
    gap_y = 0.2

    for idx, (abbr, name, tags, color) in enumerate(agents):
        r = idx // cols
        c = idx % cols
        x = Inches(grid_x + (cell_w + gap_x) * c)
        y = Inches(grid_y + (cell_h + gap_y) * r)
        w = Inches(cell_w)
        h = Inches(cell_h)
        card = add_card(
            slide, x, y, w, h,
            fill_color=CARD_BG,
            border_color=CARD_BORDER,
        )
        # 左侧彩色缩写图标
        icon_size = Inches(0.7)
        icon = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.18),
            y + Inches(0.18),
            icon_size,
            icon_size,
        )
        set_solid_fill(icon, color)
        set_no_line(icon)
        try:
            icon.adjustments[0] = 0.18
        except Exception:
            pass
        tf_icon = icon.text_frame
        tf_icon.margin_left = Emu(0)
        tf_icon.margin_right = Emu(0)
        tf_icon.margin_top = Emu(0)
        tf_icon.margin_bottom = Emu(0)
        tf_icon.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf_icon.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), abbr, 16, WHITE, bold=True, font_name=FONT_EN)

        # 文字
        add_multi_run_text(
            slide,
            x + Inches(1.0),
            y + Inches(0.25),
            Inches(cell_w - 1.1),
            Inches(0.5),
            [
                {"text": name, "size": 13, "color": WHITE, "bold": True, "font_name": FONT_EN},
            ],
        )
        add_text(
            slide,
            x + Inches(1.0),
            y + Inches(0.7),
            Inches(cell_w - 1.1),
            Inches(0.5),
            tags,
            size=11,
            color=TEXT_SECONDARY,
        )

    # 底部提示框
    tip = add_card(
        slide, Inches(0.6), Inches(6.35), Inches(12.13), Inches(0.75),
        fill_color=CARD_BG,
        border_color=CARD_BORDER,
    )
    add_card_text(
        tip,
        [[
            {"text": "双重定义：", "size": 12, "color": CYAN, "bold": True},
            {"text": "每个 Agent 在 references/agent-catalog.md（叙述版）和 references/routing-rules.json（机器可读版）中同步声明 Constraints 与 Evidence Requirements。",
             "size": 12, "color": TEXT_SECONDARY},
        ]],
        anchor=MSO_ANCHOR.MIDDLE,
        padding=0.2,
    )


def build_slide_6(prs):
    """能力矩阵。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "CAPABILITY · 能力矩阵", 6)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "14 维度 ", "size": 36, "color": WHITE, "bold": True, "font_name": FONT_EN},
            {"text": "能力对比", "size": 36, "color": PURPLE, "bold": True},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "本项目 vs 普通多专家提示词的结构性差异",
    )

    # 表格 3 列 × 6 行
    table_x = Inches(0.6)
    table_y = Inches(2.55)
    table_w = Inches(12.13)
    col_widths = [2.4, 5.5, 4.23]
    row_h = Inches(0.62)
    header_h = Inches(0.5)

    rows_data = [
        ("任务路由", "选择主负责人、协同者、治理轨道", "平铺多个角色视角"),
        ("恢复能力", "状态优先恢复、resume、恢复锚点", "依赖上下文记忆"),
        ("迭代能力", "有边界的多轮优化、基线、回滚决策", "无限'再来一轮'"),
        ("发布治理", "release gate、completion evidence、修复入口", "'建议发/不发'或只看 benchmark"),
        ("上线后闭环", "post-release feedback loop", "很少覆盖上线后反馈回写"),
        ("Subagent runtime", "受控计划 + 宿主 spawn/wait/merge 证据", "角色扮演误称真实多 Agent"),
    ]

    # 表头
    headers = [
        ("维度", WHITE),
        ("本项目 ✦", GREEN),
        ("普通方案 ✕", TEXT_MUTED),
    ]
    x_cursor = table_x
    for col_idx, (htext, hcolor) in enumerate(headers):
        cell = add_card(
            slide, x_cursor, table_y, Inches(col_widths[col_idx]), header_h,
            fill_color=RGBColor(0x1A, 0x15, 0x35),
            border_color=CARD_BORDER,
        )
        tf = cell.text_frame
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_run(p.add_run(), htext, 13, hcolor, bold=True)
        x_cursor += Inches(col_widths[col_idx])

    # 数据行
    for r_idx, row in enumerate(rows_data):
        y = table_y + header_h + row_h * r_idx
        x_cursor = table_x
        for c_idx, cell_text in enumerate(row):
            if c_idx == 0:
                color = WHITE
                bold = True
            elif c_idx == 1:
                color = GREEN
                bold = False
            else:
                color = TEXT_MUTED
                bold = False
            cell = add_card(
                slide, x_cursor, y, Inches(col_widths[c_idx]), row_h,
                fill_color=CARD_BG,
                border_color=CARD_BORDER,
            )
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.18)
            tf.margin_right = Inches(0.18)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            _set_run(p.add_run(), cell_text, 11, color, bold=bold)
            x_cursor += Inches(col_widths[c_idx])

    # 底部小字
    add_text(
        slide,
        Inches(0.6),
        Inches(7.05),
        Inches(12.13),
        Inches(0.3),
        "完整 14 维度对比见 matrix.html",
        size=11,
        color=TEXT_MUTED,
    )


def build_slide_7(prs):
    """路由示例。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "ROUTING · 路由示例", 7)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "5 个典型请求 ", "size": 36, "color": WHITE, "bold": True},
            {"text": "路由结果", "size": 36, "color": PURPLE, "bold": True},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "基于任务形态、风险、技术栈自动路由到最小可辩护 workflow bundle",
    )

    routes = [
        ("微服务架构拆分规划", "Multi-Expert", "跨架构、产品、交付多领域", PURPLE),
        ("帮我重构这段 Python 代码", "Quick Slice", "Technical Trinity · 代码编辑需交付证据", CYAN),
        ("设计 Kafka 实时数据管道", "Expert Routing", "Data Pipeline Guardian · 明确领域", BLUE),
        ("发布这个版本到生产环境", "Full Workflow", "release gate + ship/hold 决策", AMBER),
        ("我想做一个用户画像功能", "Intent Confirm", "模糊想法先确认 5 个切入方向", GREEN),
    ]

    row_y = Inches(2.6)
    row_h = Inches(0.78)
    gap = Inches(0.12)
    request_w = Inches(4.0)
    arrow_w = Inches(0.6)
    result_w = Inches(2.6)
    reason_w = Inches(4.93)

    for idx, (req, res, reason, color) in enumerate(routes):
        y = row_y + (row_h + gap) * idx

        # 请求卡片
        req_card = add_card(
            slide, Inches(0.6), y, request_w, row_h,
            fill_color=CARD_BG,
            border_color=CARD_BORDER,
        )
        tf = req_card.text_frame
        tf.margin_left = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _set_run(p.add_run(), req, 13, WHITE, bold=True)

        # 箭头
        add_text(
            slide,
            Inches(0.6) + request_w,
            y,
            arrow_w,
            row_h,
            "→",
            size=22,
            color=TEXT_MUTED,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            font_name=FONT_EN,
        )

        # 结果卡片（彩色边框）
        res_x = Inches(0.6) + request_w + arrow_w
        res_card = add_card(
            slide, res_x, y, result_w, row_h,
            fill_color=CARD_BG,
            border_color=color,
            border_width=1.25,
        )
        tf = res_card.text_frame
        tf.margin_left = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), res, 13, color, bold=True, font_name=FONT_EN)

        # 原因
        reason_x = res_x + result_w
        reason_card = add_card(
            slide, reason_x, y, reason_w, row_h,
            fill_color=CARD_BG,
            border_color=CARD_BORDER,
        )
        tf = reason_card.text_frame
        tf.margin_left = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _set_run(p.add_run(), reason, 12, TEXT_SECONDARY)


def build_slide_8(prs):
    """Quick Start。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "QUICK START · 快速上手", 8)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "3 种", "size": 36, "color": WHITE, "bold": True},
            {"text": "调用方式", "size": 36, "color": PURPLE, "bold": True},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "默认手动模式，显式 /auto 才进入自动运行",
    )

    cards = [
        ("① 手动模式", CYAN, "$virtual-intelligent-dev-team 帮我评估这次重构的最佳负责人和执行顺序。", "高风险任务逐轮确认，输出 lead route + workflow bundle + resume anchor"),
        ("② /auto setup", PURPLE, "$virtual-intelligent-dev-team /auto setup 这个项目级迁移。", "建立 automation state + 恢复锚点，为后续 go 或 resume 准备状态"),
        ("③ /auto go", GREEN, "$virtual-intelligent-dev-team /auto go → 进入自动执行", "有边界执行 + 状态优先恢复 + 保留 resume 锚点"),
    ]

    card_y = Inches(2.6)
    card_h = Inches(3.4)
    card_w = Inches(3.85)
    gap = Inches(0.25)
    start_x = Inches(0.6)

    for idx, (title, color, code, desc) in enumerate(cards):
        x = start_x + (card_w + gap) * idx
        card = add_card(
            slide, x, card_y, card_w, card_h,
            fill_color=CARD_BG,
            border_color=color,
            border_width=1.5,
        )
        # 标题
        add_text(
            slide,
            x + Inches(0.2),
            card_y + Inches(0.2),
            card_w - Inches(0.4),
            Inches(0.5),
            title,
            size=18,
            color=color,
            bold=True,
        )
        # 代码块背景
        code_bg = add_card(
            slide,
            x + Inches(0.2),
            card_y + Inches(0.85),
            card_w - Inches(0.4),
            Inches(0.9),
            fill_color=RGBColor(0x0F, 0x08, 0x20),
            border_color=CARD_BORDER,
            shape_type=MSO_SHAPE.RECTANGLE,
        )
        tf = code_bg.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _set_run(p.add_run(), code, 11, GREEN, bold=False, font_name="Courier New")

        # 描述
        add_text(
            slide,
            x + Inches(0.2),
            card_y + Inches(1.95),
            card_w - Inches(0.4),
            Inches(1.3),
            desc,
            size=12,
            color=TEXT_SECONDARY,
        )

    # 底部提示框（琥珀色边框）
    tip = add_card(
        slide, Inches(0.6), Inches(6.25), Inches(12.13), Inches(0.8),
        fill_color=CARD_BG,
        border_color=AMBER,
        border_width=1.25,
    )
    add_card_text(
        tip,
        [[
            {"text": "设计原则：", "size": 13, "color": AMBER, "bold": True},
            {"text": "自动化能力越强，越不能默认打开。默认 manual，只有显式 /auto 才进入自动运行，保持 setup → go 两阶段协议。",
             "size": 13, "color": WHITE},
        ]],
        anchor=MSO_ANCHOR.MIDDLE,
        padding=0.2,
    )


def build_slide_9(prs):
    """生态集成。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "ECOSYSTEM · 生态集成", 9)

    add_section_title(
        slide,
        Inches(0.95),
        [
            {"text": "完整", "size": 36, "color": WHITE, "bold": True},
            {"text": "生态覆盖", "size": 36, "color": PURPLE, "bold": True},
        ],
        size=36,
    )
    add_subtitle(
        slide,
        Inches(1.85),
        "从语言路由到工作流包到阶段专家团的完整集成",
    )

    # 顶部 4 个 KPI 卡片
    kpis = [
        ("13", "Language Profiles"),
        ("7", "Workflow Bundles"),
        ("8", "Specialist Agents"),
        ("7", "Closure Layers"),
    ]
    kpi_y = Inches(2.55)
    kpi_h = Inches(1.4)
    total_w = 12.13
    gap = 0.2
    n = 4
    kpi_w = (total_w - gap * (n - 1)) / n
    start_x = 0.6

    for idx, (num, label) in enumerate(kpis):
        x = Inches(start_x + (kpi_w + gap) * idx)
        w = Inches(kpi_w)
        card = add_card(
            slide, x, kpi_y, w, kpi_h,
            fill_color=CARD_BG,
            border_color=PURPLE,
            border_width=1.0,
        )
        tf = card.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), num, 40, PURPLE, bold=True, font_name=FONT_EN)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _set_run(p2.add_run(), label, 12, TEXT_SECONDARY, bold=False, font_name=FONT_EN)

    # 下方 2 个说明卡片
    desc_y = Inches(4.25)
    desc_h = Inches(2.55)
    desc_w = Inches(5.95)
    gap_x = Inches(0.25)

    # 卡片1：三层语言解耦
    card1 = add_card(
        slide, Inches(0.6), desc_y, desc_w, desc_h,
        fill_color=CARD_BG,
        border_color=BLUE,
        border_width=1.25,
    )
    add_card_text(
        card1,
        [
            [{"text": "三层语言解耦", "size": 16, "color": BLUE, "bold": True}],
            [{"text": "Routing（JSON）决定 lead agent", "size": 12, "color": TEXT_SECONDARY}],
            [{"text": "Context（YAML）注入生态默认值", "size": 12, "color": TEXT_SECONDARY}],
            [{"text": "Constraints（YAML + Manifest）注入语言护栏", "size": 12, "color": TEXT_SECONDARY}],
            [{"text": "覆盖 python/go/nodejs/rust/java/kotlin/swift/cpp/csharp/php/ruby/elixir/scala",
             "size": 11, "color": TEXT_MUTED}],
        ],
        padding=0.25,
    )

    # 卡片2：Stage Council Overlay
    card2_x = Inches(0.6) + desc_w + gap_x
    card2 = add_card(
        slide, card2_x, desc_y, desc_w, desc_h,
        fill_color=CARD_BG,
        border_color=PURPLE,
        border_width=1.25,
    )
    add_card_text(
        card2,
        [
            [{"text": "Stage Council Overlay", "size": 16, "color": PURPLE, "bold": True, "font_name": FONT_EN}],
            [{"text": "product-discovery-council：", "size": 12, "color": WHITE, "bold": True},
             {"text": "产品战略、PRD、用户研究、竞品、路线图", "size": 12, "color": TEXT_SECONDARY}],
            [{"text": "prototype-design-council：", "size": 12, "color": WHITE, "bold": True},
             {"text": "高保真原型、设计系统、可访问性审查", "size": 12, "color": TEXT_SECONDARY}],
            [{"text": "在 World-Class Product Architect 下按需展开，不替换顶层 lead",
             "size": 11, "color": TEXT_MUTED}],
        ],
        padding=0.25,
    )


def build_slide_10(prs):
    """CTA。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_eyebrow(slide, "GET STARTED · 开始使用", 10)

    # 中央大卡片（紫色渐变背景用纯色 PURPLE 模拟）
    card_x = Inches(0.9)
    card_y = Inches(1.3)
    card_w = Inches(11.53)
    card_h = Inches(5.5)

    main_card = add_card(
        slide, card_x, card_y, card_w, card_h,
        fill_color=RGBColor(0x1A, 0x15, 0x35),
        border_color=PURPLE,
        border_width=2.0,
    )

    # 标题
    add_text(
        slide,
        card_x + Inches(0.5),
        card_y + Inches(0.4),
        card_w - Inches(1.0),
        Inches(0.8),
        "开始你的第一个闭环任务",
        size=40,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 描述
    add_paragraphs(
        slide,
        card_x + Inches(1.5),
        card_y + Inches(1.5),
        card_w - Inches(3.0),
        Inches(1.0),
        [
            [
                {"text": "适合接手'单个专家已经不够、单轮回答也不够'的复杂软件任务", "size": 16, "color": TEXT_SECONDARY},
            ],
            [
                {"text": "从手动模式开始，按需进入 /auto 自动运行", "size": 16, "color": TEXT_SECONDARY},
            ],
        ],
        align=PP_ALIGN.CENTER,
    )

    # 链接行（5 个按钮）
    btn_y = card_y + Inches(3.0)
    btn_h = Inches(0.5)
    btn_specs = [
        ("★ GitHub Repository", PURPLE, WHITE, True),
        ("返回首页", CARD_BG, TEXT_SECONDARY, False),
        ("七层闭环架构", CARD_BG, TEXT_SECONDARY, False),
        ("8 个 Agent 角色", CARD_BG, TEXT_SECONDARY, False),
        ("能力矩阵对比", CARD_BG, TEXT_SECONDARY, False),
    ]

    # 计算按钮宽度与间距
    total_btn_w = 0
    widths = []
    for text, _, _, _ in btn_specs:
        w = 0.16 * len(text) + 0.5
        widths.append(w)
        total_btn_w += w
    total_btn_w += 0.15 * (len(btn_specs) - 1)

    start_x = (13.333 - total_btn_w) / 2
    cur_x = start_x
    for idx, (text, bg, tc, bold) in enumerate(btn_specs):
        w = widths[idx]
        btn = add_card(
            slide, Inches(cur_x), btn_y, Inches(w), btn_h,
            fill_color=bg,
            border_color=PURPLE if bold else CARD_BORDER,
            border_width=1.5 if bold else 1.0,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        try:
            btn.adjustments[0] = 0.3
        except Exception:
            pass
        tf = btn.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), text, 12, tc, bold=bold)
        cur_x += w + 0.15

    # 底部小字
    add_paragraphs(
        slide,
        card_x + Inches(1.5),
        card_y + Inches(4.0),
        card_w - Inches(3.0),
        Inches(1.2),
        [
            [
                {"text": "Virtual Intelligent Dev Team · v5.7.6 · MIT License", "size": 13, "color": TEXT_MUTED},
            ],
            [
                {"text": "github.com/fxbin/virtual-intelligent-dev-team", "size": 13, "color": CYAN, "font_name": FONT_EN},
            ],
        ],
        align=PP_ALIGN.CENTER,
    )


# ============================================================
# 主流程
# ============================================================
def main():
    """构造 Presentation 并依次生成 10 页幻灯片。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        build_slide_1, build_slide_2, build_slide_3, build_slide_4,
        build_slide_5, build_slide_6, build_slide_7, build_slide_8,
        build_slide_9, build_slide_10,
    ]

    for idx, builder in enumerate(builders, start=1):
        builder(prs)
        print(f"[ok] slide {idx:02d} built")

    out_path = "/Users/fxbin/Desktop/Project/AIProject/skill-hub/virtual-intelligent-dev-team/docs/assets/intro-deck/virtual-intelligent-dev-team-intro.pptx"
    prs.save(out_path)
    print(f"\n[saved] {out_path}")
    return out_path


if __name__ == "__main__":
    main()
